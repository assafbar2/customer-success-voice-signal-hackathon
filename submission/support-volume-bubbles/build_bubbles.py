#!/usr/bin/env python3
"""Build an editable PowerPoint bubble chart of support-adjacent volumes.

Bubble *area* is proportional to volume (radius ∝ √value). Support is
highlighted with bold type, an accent fill, a ring, and a callout chip.
"""

from __future__ import annotations

import math
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUT = Path(__file__).resolve().parent / "support-volume-bubbles.pptx"

CHANNELS: list[dict] = [
    {"label": "Documentation", "value": 300_000, "display": "300,000"},
    {"label": "Help centers total articles views", "value": 152_000, "display": "152,000"},
    {"label": "Kapa", "value": 11_000, "display": "11,000"},
    {"label": "GitHub issues & discussions", "value": 2_700, "display": "2,700"},
    {"label": "Support", "value": 2_500, "display": "2,500", "highlight": True},
    {"label": "Social channels", "value": 600, "display": "600"},
    {"label": "Discord & community", "value": 200, "display": "200"},
    {"label": "Feedback systems", "value": 200, "display": "200"},
]

FOOTNOTE = "We redesigned one of the smallest, clearest surfaces first."

BG = RGBColor(0xF7, 0xF5, 0xF2)
INK = RGBColor(0x1A, 0x17, 0x14)
INK_MUTED = RGBColor(0x5C, 0x55, 0x4C)
BUBBLE_FILL = RGBColor(0xD8, 0xCF, 0xC3)
BUBBLE_LINE = RGBColor(0xB9, 0xAE, 0xA0)
SUPPORT_FILL = RGBColor(0xE8, 0xA5, 0x4B)
SUPPORT_RING = RGBColor(0x8A, 0x3B, 0x12)
SUPPORT_INK = RGBColor(0x2A, 0x14, 0x08)
CHIP_INK = RGBColor(0xFF, 0xFB, 0xF6)

LABEL_PT = 12
VALUE_PT = 14


@dataclass
class Bubble:
    channel: dict
    cx: float
    cy: float
    r: float


def set_run(paragraph, text: str, *, size_pt: float, bold: bool, color: RGBColor) -> None:
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def radius_for(value: float, max_value: float, max_r: float) -> float:
    return max_r * math.sqrt(value / max_value)


def pack_bubbles(channels: list[dict], origin_x: float, origin_y: float, width: float, height: float) -> list[Bubble]:
    """Largest-first greedy circle packing inside the frame."""
    max_v = max(c["value"] for c in channels)
    max_r = min(width, height) * 0.40
    ordered = sorted(channels, key=lambda c: c["value"], reverse=True)

    placed: list[Bubble] = []
    cx0 = origin_x + width * 0.36
    cy0 = origin_y + height * 0.50
    gap = float(Emu(35_000))

    for i, ch in enumerate(ordered):
        r = radius_for(ch["value"], max_v, max_r)
        if i == 0:
            placed.append(Bubble(ch, cx0, cy0, r))
            continue

        best = None
        best_score = float("inf")
        for host in placed:
            steps = 48
            for ang_i in range(steps):
                ang = (ang_i / steps) * 2 * math.pi
                dist = host.r + r + gap
                cx = host.cx + math.cos(ang) * dist
                cy = host.cy + math.sin(ang) * dist
                if cx - r < origin_x or cy - r < origin_y:
                    continue
                if cx + r > origin_x + width or cy + r > origin_y + height:
                    continue
                if any(math.hypot(cx - p.cx, cy - p.cy) < p.r + r + gap * 0.6 for p in placed):
                    continue
                score = math.hypot(cx - cx0, cy - cy0)
                if ch.get("highlight"):
                    # Prefer a clear pocket on the lower-right for the callout.
                    score += (cy0 - cy) * 0.25
                    score -= (cx - cx0) * 0.05
                if score < best_score:
                    best_score = score
                    best = (cx, cy)

        if best is None:
            cx = origin_x + width * 0.78
            cy = origin_y + height * 0.78
            best = (cx, cy)
        placed.append(Bubble(ch, best[0], best[1], r))
    return placed


def add_oval(slide, cx: float, cy: float, r: float, *, fill: RGBColor, line: RGBColor, line_pt: float):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        int(cx - r),
        int(cy - r),
        int(r * 2),
        int(r * 2),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(line_pt)
    return shape


def add_bubble(slide, bubble: Bubble) -> None:
    ch = bubble.channel
    highlight = bool(ch.get("highlight"))
    fill = SUPPORT_FILL if highlight else BUBBLE_FILL
    line = SUPPORT_RING if highlight else BUBBLE_LINE
    line_pt = 2.75 if highlight else 1.0
    shape = add_oval(slide, bubble.cx, bubble.cy, bubble.r, fill=fill, line=line, line_pt=line_pt)

    ink = SUPPORT_INK if highlight else INK
    bold = highlight
    size = bubble.r * 2
    min_inner = float(Emu(950_000))

    if size >= min_inner:
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(2)
        set_run(p, ch["label"], size_pt=LABEL_PT, bold=bold, color=ink)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        set_run(p2, ch["display"], size_pt=VALUE_PT, bold=bold, color=ink)
        return

    label_w = float(Inches(1.6))
    label_h = float(Inches(0.55))
    lx = bubble.cx + bubble.r + float(Emu(50_000))
    ly = bubble.cy - label_h / 2
    if lx + label_w > float(Inches(12.9)):
        lx = bubble.cx - bubble.r - label_w - float(Emu(40_000))
    box = slide.shapes.add_textbox(int(lx), int(ly), int(label_w), int(label_h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    set_run(p, ch["label"], size_pt=LABEL_PT, bold=bold, color=ink)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    set_run(p2, ch["display"], size_pt=VALUE_PT, bold=bold, color=ink)


def add_support_highlight(slide, bubble: Bubble) -> None:
    """Halo ring + caption chip so Support numbers can't be missed."""
    pad = float(Emu(110_000))
    halo = add_oval(
        slide,
        bubble.cx,
        bubble.cy,
        bubble.r + pad,
        fill=BG,
        line=SUPPORT_RING,
        line_pt=1.75,
    )
    halo.fill.background()

    chip_w = float(Inches(1.85))
    chip_h = float(Inches(0.38))
    chip = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        int(bubble.cx - chip_w / 2),
        int(bubble.cy + bubble.r + float(Emu(80_000))),
        int(chip_w),
        int(chip_h),
    )
    chip.adjustments[0] = 0.5
    chip.fill.solid()
    chip.fill.fore_color.rgb = SUPPORT_RING
    chip.line.fill.background()
    tf = chip.text_frame
    tf.clear()
    tf.anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p, "Support   2,500", size_pt=13, bold=True, color=CHIP_INK)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(12.2), Inches(0.4))
    set_run(title.text_frame.paragraphs[0], "Where attention already lives", size_pt=22, bold=False, color=INK)

    subtitle = slide.shapes.add_textbox(Inches(0.55), Inches(0.68), Inches(12.2), Inches(0.32))
    set_run(
        subtitle.text_frame.paragraphs[0],
        "Bubble area = volume. Same type everywhere — Support is bold, ringed, and called out.",
        size_pt=12,
        bold=False,
        color=INK_MUTED,
    )

    bubbles = pack_bubbles(
        CHANNELS,
        float(Inches(0.45)),
        float(Inches(1.15)),
        float(Inches(12.4)),
        float(Inches(4.85)),
    )

    for b in sorted(bubbles, key=lambda b: b.r, reverse=True):
        if b.channel.get("highlight"):
            continue
        add_bubble(slide, b)

    support = next(b for b in bubbles if b.channel.get("highlight"))
    add_support_highlight(slide, support)
    add_bubble(slide, support)

    note = slide.shapes.add_textbox(Inches(0.55), Inches(6.2), Inches(12.2), Inches(0.28))
    set_run(
        note.text_frame.paragraphs[0],
        "Area scales with volume (Documentation ≈ 1500× Feedback / Discord). GitHub = 2,700.",
        size_pt=11,
        bold=False,
        color=INK_MUTED,
    )

    foot = slide.shapes.add_textbox(Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.4))
    set_run(foot.text_frame.paragraphs[0], FOOTNOTE, size_pt=14, bold=False, color=INK)

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    print(f"Wrote {build()}")
