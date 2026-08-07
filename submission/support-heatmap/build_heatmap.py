#!/usr/bin/env python3
"""Build an editable PowerPoint treemap of support-adjacent surfaces.

Tile *area* is proportional to volume. All labels share one font/size;
only "Support" is bold. Native shapes + text for PowerPoint editing.
"""

from __future__ import annotations

import math
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "support-channel-heatmap.pptx"

CHANNELS: list[dict] = [
    {"label": "Documentation", "value": 300_000, "display": "300,000"},
    {"label": "Help centers total articles views", "value": 152_000, "display": "152,000"},
    {"label": "Kapa", "value": 11_000, "display": "11,000"},
    {"label": "GitHub issues & discussions", "value": 2_500, "display": "2500"},
    {"label": "Support", "value": 2_500, "display": "2500", "bold": True},
    {"label": "Social channels", "value": 600, "display": "600"},
    {"label": "Feedback systems", "value": 200, "display": "200"},
]

DISCORD = {"label": "Discord & community", "value": None, "display": "[TBD]"}
FOOTNOTE = "We redesigned one of the smallest, clearest surfaces first."

BG = RGBColor(0xF7, 0xF5, 0xF2)
INK = RGBColor(0x1A, 0x17, 0x14)
INK_MUTED = RGBColor(0x5C, 0x55, 0x4C)
TILE_UNKNOWN = RGBColor(0xE4, 0xE0, 0xDA)
GAP = Inches(0.012)

LABEL_PT = 12
VALUE_PT = 15


def lerp(a: float, b: float, t: float) -> float:
    return a + (b - a) * t


def heat_color(t: float) -> RGBColor:
    t = max(0.0, min(1.0, t))
    stops = [
        (0.0, (0xF0, 0xE6, 0xD8)),
        (0.35, (0xF2, 0xC4, 0x7A)),
        (0.7, (0xE8, 0x8B, 0x4A)),
        (1.0, (0xC4, 0x45, 0x2D)),
    ]
    for i in range(len(stops) - 1):
        t0, c0 = stops[i]
        t1, c1 = stops[i + 1]
        if t <= t1:
            u = 0.0 if t1 == t0 else (t - t0) / (t1 - t0)
            return RGBColor(
                int(lerp(c0[0], c1[0], u)),
                int(lerp(c0[1], c1[1], u)),
                int(lerp(c0[2], c1[2], u)),
            )
    return RGBColor(*stops[-1][1])


def contrast_ink(fill: RGBColor) -> RGBColor:
    r, g, b = fill[0] / 255, fill[1] / 255, fill[2] / 255
    lum = 0.2126 * r + 0.7152 * g + 0.0722 * b
    return RGBColor(0xFF, 0xFB, 0xF6) if lum < 0.55 else INK


def heat_for(value: int, lo: float, hi: float) -> float:
    if hi <= lo:
        return 0.5
    return (math.log10(value) - lo) / (hi - lo)


@dataclass
class Rect:
    x: float
    y: float
    w: float
    h: float


def sized_layout(channels: list[dict], x: float, y: float, w: float, h: float) -> list[tuple[dict, Rect]]:
    """Area-proportional treemap tuned for this extreme volume range.

    1) Documentation = left vertical band (share of total width).
    2) Help centers = top of the right band (share of right height).
    3) Remaining surfaces share the bottom-right strip, widths ∝ volume.

    Area stays linear with volume; the bottom strip stays wide enough to read.
    """
    by_name = {c["label"]: c for c in channels}
    total = sum(c["value"] for c in channels)

    doc = by_name["Documentation"]
    help_ = by_name["Help centers total articles views"]
    small = [c for c in channels if c["label"] not in (doc["label"], help_["label"])]
    small = sorted(small, key=lambda c: c["value"], reverse=True)

    doc_w = w * (doc["value"] / total)
    right_x = x + doc_w
    right_w = w - doc_w

    right_total = total - doc["value"]
    help_h = h * (help_["value"] / right_total)
    small_y = y + help_h
    small_h = h - help_h
    small_total = sum(c["value"] for c in small)

    out: list[tuple[dict, Rect]] = [
        (doc, Rect(x, y, doc_w, h)),
        (help_, Rect(right_x, y, right_w, help_h)),
    ]
    cx = right_x
    for i, ch in enumerate(small):
        tw = right_w * (ch["value"] / small_total)
        # Last tile absorbs float remainder so the strip seals flush.
        if i == len(small) - 1:
            tw = right_x + right_w - cx
        out.append((ch, Rect(cx, small_y, tw, small_h)))
        cx += tw
    return out


def set_run(paragraph, text: str, *, size_pt: float, bold: bool, color: RGBColor) -> None:
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def add_tile(slide, rect: Rect, channel: dict, heat: float | None, *, gap: int) -> None:
    # Keep a hairline gutter without erasing tiny volume tiles.
    inset = min(gap // 2, max(int(rect.w) // 8, 0), max(int(rect.h) // 8, 0))
    left = int(round(rect.x)) + inset
    top = int(round(rect.y)) + inset
    width = max(int(round(rect.w)) - 2 * inset, 1)
    height = max(int(round(rect.h)) - 2 * inset, 1)

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.05
    fill = shape.fill
    fill.solid()
    rgb = TILE_UNKNOWN if heat is None else heat_color(heat)
    fill.fore_color.rgb = rgb
    shape.line.fill.background()

    ink = contrast_ink(rgb)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = None
    tf.anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(2)
    set_run(
        p,
        channel["label"],
        size_pt=LABEL_PT,
        bold=bool(channel.get("bold")),
        color=ink,
    )

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    set_run(p2, channel["display"], size_pt=VALUE_PT, bold=False, color=ink)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(12.2), Inches(0.4))
    set_run(title.text_frame.paragraphs[0], "Where attention already lives", size_pt=22, bold=False, color=INK)

    subtitle = slide.shapes.add_textbox(Inches(0.55), Inches(0.68), Inches(12.2), Inches(0.32))
    set_run(
        subtitle.text_frame.paragraphs[0],
        "Tile size = volume (linear). Color reinforces the same scale. Same type everywhere — Support is bold.",
        size_pt=12,
        bold=False,
        color=INK_MUTED,
    )

    origin_x = float(Inches(0.55))
    origin_y = float(Inches(1.15))
    region_w = float(Inches(12.2))
    region_h = float(Inches(4.85))
    discord_w = float(Inches(0.72))
    tree_w = region_w - discord_w - float(GAP)
    gap_emu = int(GAP)

    values = [c["value"] for c in CHANNELS]
    lo, hi = math.log10(min(values)), math.log10(max(values))

    for channel, rect in sized_layout(CHANNELS, origin_x, origin_y, tree_w, region_h):
        add_tile(slide, rect, channel, heat_for(channel["value"], lo, hi), gap=gap_emu)

    discord_rect = Rect(origin_x + tree_w + float(GAP), origin_y, discord_w, region_h)
    add_tile(slide, discord_rect, DISCORD, None, gap=gap_emu)

    note = slide.shapes.add_textbox(Inches(0.55), Inches(6.15), Inches(12.2), Inches(0.28))
    set_run(
        note.text_frame.paragraphs[0],
        "Area is linear with volume (Documentation ≈ 1500× Feedback). Discord is unmeasured — not on the size scale.",
        size_pt=11,
        bold=False,
        color=INK_MUTED,
    )

    foot = slide.shapes.add_textbox(Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.4))
    set_run(foot.text_frame.paragraphs[0], FOOTNOTE, size_pt=14, bold=False, color=INK)

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path}")
